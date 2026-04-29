#!/usr/bin/env python3
"""Native python-pptx diagram renderers — 14 diagram types drawn entirely with
autoshapes, connectors, and text boxes. No SmartArt, no SVG.

Usage:
    from diagram_shapes import draw_diagram
    draw_diagram(slide, 'pyramid', data, style)
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math

# ---------------------------------------------------------------------------
# Slide constants (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Default style — merged with user-supplied style before rendering
# ---------------------------------------------------------------------------
DEFAULT_STYLE = {
    'primary': RGBColor(0x1B, 0x36, 0x5D),
    'secondary': RGBColor(0x2E, 0x5C, 0x8A),
    'accent': RGBColor(0x0D, 0x73, 0x77),
    'bg': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x1A, 0x1A, 0x1A),
    'text_light': RGBColor(0x63, 0x6E, 0x72),
    'border': RGBColor(0xCC, 0xCC, 0xCC),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'font_heading': 'Arial',
    'font_body': 'Arial',
    'font_size_title': Pt(28),
    'font_size_heading': Pt(18),
    'font_size_body': Pt(13),
    'font_size_small': Pt(11),
}


def _resolve_style(user_style):
    """Return a merged style dict, falling back to DEFAULT_STYLE for missing keys."""
    if user_style is None:
        return dict(DEFAULT_STYLE)
    merged = dict(DEFAULT_STYLE)
    merged.update(user_style)
    return merged


# ======================================================================
# Helper Functions
# ======================================================================

def add_title_bar(slide, title, style):
    """Add a title with accent underline at the top of the slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = style.get('font_size_title', Pt(28))
    p.font.bold = True
    p.font.color.rgb = style.get('text', DEFAULT_STYLE['text'])
    p.font.name = style.get('font_heading', 'Arial')

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.68), Inches(2.0), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
    line.line.fill.background()


def add_rounded_card(slide, x, y, w, h, title, items, style,
                     fill_color=None, text_color=None):
    """Add a card with rounded corners, title, and bullet items."""
    fill_color = fill_color or style.get('bg', DEFAULT_STYLE['bg'])
    text_color = text_color or style.get('text', DEFAULT_STYLE['text'])

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = style.get('border', DEFAULT_STYLE['border'])
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(4)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = style.get('font_size_heading', Pt(18))
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = style.get('font_heading', 'Arial')

    for item in items:
        p = tf.add_paragraph()
        p.text = f'  \u2022 {item}'
        p.font.size = style.get('font_size_small', Pt(11))
        p.font.color.rgb = style.get('text_light', DEFAULT_STYLE['text_light'])
        p.font.name = style.get('font_body', 'Arial')
        p.space_before = Pt(3)

    return shape


def add_filled_block(slide, x, y, w, h, label, items, style,
                     fill_color, text_color=None):
    """Add a filled rectangular block with label and items."""
    text_color = text_color or DEFAULT_STYLE['white']

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(4)

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = style.get('font_size_heading', Pt(18))
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = style.get('font_heading', 'Arial')

    for item in items:
        p = tf.add_paragraph()
        p.text = f'  \u2022 {item}'
        p.font.size = style.get('font_size_small', Pt(11))
        p.font.color.rgb = text_color
        p.font.name = style.get('font_body', 'Arial')
        p.space_before = Pt(3)

    return shape


def add_footer_note(slide, text, style):
    """Add a small gray footer note."""
    txBox = slide.shapes.add_textbox(Inches(2), Inches(7.0), Inches(9), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = style.get('text_light', DEFAULT_STYLE['text_light'])
    p.font.name = style.get('font_body', 'Arial')
    p.alignment = PP_ALIGN.CENTER


def _add_connector(slide, x1, y1, x2, y2, style, connector_type=None, color=None):
    """Add a straight connector line between two points."""
    ct = connector_type or MSO_CONNECTOR_TYPE.STRAIGHT
    conn = slide.shapes.add_connector(ct, x1, y1, x2, y2)
    conn.line.color.rgb = color or style.get('accent', DEFAULT_STYLE['accent'])
    conn.line.width = Pt(2)
    return conn


def _add_textbox(slide, x, y, w, h, text, style, font_size=None,
                 color=None, bold=False, alignment=None, font_name=None):
    """Add a single-line text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size or style.get('font_size_body', Pt(13))
    p.font.color.rgb = color or style.get('text', DEFAULT_STYLE['text'])
    p.font.bold = bold
    p.font.name = font_name or style.get('font_body', 'Arial')
    if alignment is not None:
        p.alignment = alignment
    return txBox


def _add_circle(slide, cx, cy, r, style, fill_color=None, line_color=None,
                line_width=Pt(2)):
    """Add a circle (oval inscribed in a square)."""
    d = 2 * r
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r, cy - r, d, d
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = line_color or style.get('accent', DEFAULT_STYLE['accent'])
    shape.line.width = line_width
    return shape


def _set_cell_text(cell, text, font_size, color_rgb, bold=False, font_name='Arial'):
    """Helper to set text on a table cell paragraph."""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.color.rgb = color_rgb
    p.font.bold = bold
    p.font.name = font_name


# ======================================================================
# 1. Pyramid
# ======================================================================

def draw_pyramid(slide, data, style):
    """Draw a hierarchical pyramid with horizontal stacked layers.

    data = {'layers': [{'label': str, 'items': [str]}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    layers = data.get('layers', [])
    n = len(layers)
    if n == 0:
        return

    top_y = Inches(1.1)
    total_height = Inches(5.6)
    max_width = Inches(11.0)
    min_width = Inches(4.0)
    center_x = Inches(6.666)

    color_ramp = [
        style.get('accent', DEFAULT_STYLE['accent']),       # top
        style.get('secondary', DEFAULT_STYLE['secondary']),  # middle
        style.get('primary', DEFAULT_STYLE['primary']),      # bottom
    ]

    for i, layer in enumerate(layers):
        t = i / max(n - 1, 1)
        w = max_width - (max_width - min_width) * t
        fill = color_ramp[min(i, len(color_ramp) - 1)]

        layer_h = total_height / n - Pt(6)
        y = top_y + i * (total_height / n)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            center_x - w / 2, y, w, layer_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        tf.margin_top = Pt(6)

        p = tf.paragraphs[0]
        p.text = layer.get('label', '')
        p.font.size = style.get('font_size_heading', Pt(18))
        p.font.bold = True
        p.font.color.rgb = DEFAULT_STYLE['white']
        p.font.name = style.get('font_heading', 'Arial')

        for item in layer.get('items', []):
            p = tf.add_paragraph()
            p.text = f'  \u2022 {item}'
            p.font.size = style.get('font_size_small', Pt(11))
            p.font.color.rgb = DEFAULT_STYLE['white']
            p.font.name = style.get('font_body', 'Arial')
            p.space_before = Pt(2)

    add_footer_note(slide, '\u81ea\u4e0b\u800c\u4e0a', style)


# ======================================================================
# 2. Hub-and-Spoke
# ======================================================================

def draw_hub_spoke(slide, data, style):
    """Draw a hub-and-spoke diagram with center circle and branch cards.

    data = {'center': {'label': str}, 'branches': [{'label': str, 'items': [str]}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    center = data.get('center', {'label': 'Hub'})
    branches = data.get('branches', [])
    n_branches = len(branches)

    hub_cx = Inches(6.666)
    hub_cy = Inches(3.8)
    hub_r = Inches(0.65)

    _add_circle(slide, hub_cx, hub_cy, hub_r, style,
                fill_color=style.get('primary', DEFAULT_STYLE['primary']),
                line_color=style.get('primary', DEFAULT_STYLE['primary']))
    _add_textbox(slide, hub_cx - Inches(1), hub_cy - Pt(12),
                 Inches(2), Inches(0.5),
                 center.get('label', ''), style,
                 font_size=style.get('font_size_body', Pt(13)),
                 color=DEFAULT_STYLE['white'], bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_name=style.get('font_heading', 'Arial'))

    card_w = Inches(3.5)
    card_h = Inches(2.0)

    if n_branches <= 2:
        positions = [
            (Inches(0.6), Inches(2.8)),
            (Inches(8.8), Inches(2.8)),
        ]
    elif n_branches <= 4:
        positions = [
            (Inches(0.6), Inches(1.5)),
            (Inches(8.8), Inches(1.5)),
            (Inches(0.6), Inches(4.6)),
            (Inches(8.8), Inches(4.6)),
        ]
    else:
        positions = [
            (Inches(0.6), Inches(1.2)),
            (Inches(8.8), Inches(1.2)),
            (Inches(0.6), Inches(2.9)),
            (Inches(8.8), Inches(2.9)),
            (Inches(0.6), Inches(4.6)),
            (Inches(8.8), Inches(4.6)),
        ]

    for idx, branch in enumerate(branches):
        if idx >= len(positions):
            break
        bx, by = positions[idx]

        add_rounded_card(slide, bx, by, card_w, card_h,
                         branch.get('label', ''), branch.get('items', []),
                         style)

        cx = bx + card_w / 2
        cy = by + card_h / 2
        dx = hub_cx - cx
        dy = hub_cy - cy
        dist = math.sqrt(float(dx) ** 2 + float(dy) ** 2)
        if dist > 0:
            ratio = float(hub_r) / dist
            sx = cx + dx * ratio
            sy = cy + dy * ratio
        else:
            sx, sy = cx, cy

        _add_connector(slide, sx, sy, cx, cy, style)

    add_footer_note(slide, '\u4ee5\u4e2d\u5fc3\u4e3a\u6838\u5fc3\uff0c\u8f90\u5c04\u5404\u4e2a\u5206\u652f', style)


# ======================================================================
# 3. Dual-Gears
# ======================================================================

def draw_dual_gears(slide, data, style):
    """Draw two interlocking gears side by side.

    data = {'left': {'label': str, 'items': [str]},
            'right': {'label': str, 'items': [str]},
            'arrow_label': str}
    """
    add_title_bar(slide, data.get('title', ''), style)

    left = data.get('left', {'label': '', 'items': []})
    right = data.get('right', {'label': '', 'items': []})
    arrow_label = data.get('arrow_label', '')

    left_cx = Inches(4.5)
    right_cx = Inches(8.8)
    gear_cy = Inches(3.6)
    gear_r = Inches(1.2)

    left_fill = style.get('secondary', DEFAULT_STYLE['secondary'])
    left_fill_light = RGBColor(
        min(left_fill[0] + 80, 255),
        min(left_fill[1] + 80, 255),
        min(left_fill[2] + 80, 255)
    )
    _add_circle(slide, left_cx, gear_cy, gear_r, style,
                fill_color=left_fill_light,
                line_color=left_fill,
                line_width=Pt(3))

    _add_textbox(slide, left_cx - Inches(1), gear_cy - Pt(14),
                 Inches(2), Inches(0.4),
                 left.get('label', ''), style,
                 font_size=style.get('font_size_heading', Pt(18)),
                 color=style.get('text', DEFAULT_STYLE['text']),
                 bold=True, alignment=PP_ALIGN.CENTER)

    right_fill = style.get('accent', DEFAULT_STYLE['accent'])
    right_fill_light = RGBColor(
        min(right_fill[0] + 80, 255),
        min(right_fill[1] + 80, 255),
        min(right_fill[2] + 80, 255)
    )
    _add_circle(slide, right_cx, gear_cy, gear_r, style,
                fill_color=right_fill_light,
                line_color=right_fill,
                line_width=Pt(3))

    _add_textbox(slide, right_cx - Inches(1), gear_cy - Pt(14),
                 Inches(2), Inches(0.4),
                 right.get('label', ''), style,
                 font_size=style.get('font_size_heading', Pt(18)),
                 color=style.get('text', DEFAULT_STYLE['text']),
                 bold=True, alignment=PP_ALIGN.CENTER)

    item_y = Inches(2.0)
    for item in left.get('items', []):
        _add_textbox(slide, Inches(1.2), item_y, Inches(2.8), Inches(0.3),
                     f'\u2022 {item}', style,
                     font_size=style.get('font_size_small', Pt(11)),
                     color=style.get('text_light', DEFAULT_STYLE['text_light']))
        item_y += Inches(0.35)

    item_y = Inches(2.0)
    for item in right.get('items', []):
        _add_textbox(slide, Inches(9.3), item_y, Inches(2.8), Inches(0.3),
                     f'\u2022 {item}', style,
                     font_size=style.get('font_size_small', Pt(11)),
                     color=style.get('text_light', DEFAULT_STYLE['text_light']))
        item_y += Inches(0.35)

    arrow_y = gear_cy - gear_r - Inches(0.8)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.UP_ARROW, Inches(6.1), arrow_y, Inches(1.1), Inches(0.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
    arrow.line.fill.background()

    _add_textbox(slide, Inches(5.3), arrow_y + Inches(0.6),
                 Inches(2.8), Inches(0.3),
                 arrow_label, style,
                 font_size=style.get('font_size_small', Pt(11)),
                 color=style.get('accent', DEFAULT_STYLE['accent']),
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_footer_note(slide, '\u53cc\u9f7f\u8054\u52a8\uff0c\u76f8\u4e92\u9a71\u52a8', style)


# ======================================================================
# 4. Tension-Triangle
# ======================================================================

def draw_tension_triangle(slide, data, style):
    """Draw a triangular tension diagram with 3 nodes, edge labels, and center icon.

    data = {'nodes': [{'label': str}, ...],
            'edges': [{'from': int, 'to': int, 'label': str}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    nodes = data.get('nodes', [])
    edges = data.get('edges', [])

    if len(nodes) < 3:
        return

    top = (Inches(6.666), Inches(1.5))
    bottom_left = (Inches(2.3), Inches(6.0))
    bottom_right = (Inches(11.0), Inches(6.0))
    positions_schema = [top, bottom_left, bottom_right]
    node_r = Inches(0.45)

    edge_color = style.get('primary', DEFAULT_STYLE['primary'])
    for edge in edges:
        f_idx = edge.get('from', 0)
        t_idx = edge.get('to', 0)
        if f_idx < 3 and t_idx < 3:
            x1, y1 = positions_schema[f_idx]
            x2, y2 = positions_schema[t_idx]

            dx = float(x2) - float(x1)
            dy = float(y2) - float(y1)
            dist = math.sqrt(dx**2 + dy**2)
            if dist > 0:
                ratio = float(node_r + Pt(2)) / dist
                sx = x1 + dx * ratio
                sy = y1 + dy * ratio
                ex = x2 - dx * ratio
                ey = y2 - dy * ratio
            else:
                sx, sy, ex, ey = x1, y1, x2, y2

            _add_connector(slide, sx, sy, ex, ey, style, color=edge_color)

            mx = (x1 + x2) / 2
            my = (y1 + y2) / 2
            label = edge.get('label', '')
            if label:
                _add_textbox(slide, mx - Inches(1.2), my - Inches(0.3),
                             Inches(2.4), Inches(0.4),
                             label, style,
                             font_size=style.get('font_size_small', Pt(11)),
                             color=edge_color,
                             bold=True, alignment=PP_ALIGN.CENTER)

    for idx, node in enumerate(nodes):
        if idx >= 3:
            break
        nx, ny = positions_schema[idx]

        _add_circle(slide, nx, ny, node_r, style,
                    fill_color=style.get('accent', DEFAULT_STYLE['accent']),
                    line_color=style.get('accent', DEFAULT_STYLE['accent']))

        _add_textbox(slide, nx - Inches(1.2), ny + node_r + Pt(6),
                     Inches(2.4), Inches(0.4),
                     node.get('label', ''), style,
                     font_size=style.get('font_size_body', Pt(13)),
                     color=style.get('text', DEFAULT_STYLE['text']),
                     bold=True, alignment=PP_ALIGN.CENTER)

    center_cx = Inches(6.666)
    center_cy = Inches(4.5)

    center_shape = _add_circle(slide, center_cx, center_cy, Inches(0.3), style,
                               fill_color=style.get('secondary', DEFAULT_STYLE['secondary']),
                               line_color=style.get('secondary', DEFAULT_STYLE['secondary']))
    _add_textbox(slide, center_cx - Inches(0.6), center_cy - Pt(12),
                 Inches(1.2), Inches(0.4),
                 '\u25b2', style,
                 font_size=Pt(16),
                 color=DEFAULT_STYLE['white'],
                 alignment=PP_ALIGN.CENTER)


# ======================================================================
# 5. Bubble-Matrix
# ======================================================================

def draw_bubble_matrix(slide, data, style):
    """Draw a 2x2 quadrant matrix with bubbles.

    data = {'x_axis': str, 'y_axis': str,
            'quadrants': {...},
            'bubbles': [{...}],
            'takeaway': str}
    """
    add_title_bar(slide, data.get('title', ''), style)

    x_axis = data.get('x_axis', '')
    y_axis = data.get('y_axis', '')
    bubbles = data.get('bubbles', [])
    takeaway = data.get('takeaway', '')
    quadrants = data.get('quadrants', {})

    grid_x = Inches(2.5)
    grid_y = Inches(1.3)
    grid_w = Inches(9.0)
    grid_h = Inches(5.0)

    grid_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, grid_x, grid_y, grid_w, grid_h
    )
    grid_shape.fill.background()
    grid_shape.line.color.rgb = style.get('border', DEFAULT_STYLE['border'])
    grid_shape.line.width = Pt(1)

    mid_x = grid_x + grid_w / 2
    mid_y = grid_y + grid_h / 2

    _add_connector(slide, grid_x, mid_y, grid_x + grid_w, mid_y, style,
                   color=style.get('border', DEFAULT_STYLE['border']))
    _add_connector(slide, mid_x, grid_y, mid_x, grid_y + grid_h, style,
                   color=style.get('border', DEFAULT_STYLE['border']))

    _add_textbox(slide, grid_x + grid_w / 2 - Inches(1.5), grid_y + grid_h + Pt(8),
                 Inches(3), Inches(0.3),
                 x_axis, style,
                 font_size=style.get('font_size_small', Pt(11)),
                 color=style.get('text_light', DEFAULT_STYLE['text_light']),
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, grid_x - Inches(1.8), grid_y + grid_h / 2 - Inches(0.2),
                 Inches(1.6), Inches(0.4),
                 y_axis, style,
                 font_size=style.get('font_size_small', Pt(11)),
                 color=style.get('text_light', DEFAULT_STYLE['text_light']),
                 alignment=PP_ALIGN.CENTER)

    q_data = quadrants if quadrants else {}
    q_top_left = q_data.get('top_left', q_data.get('q1', ''))
    q_top_right = q_data.get('top_right', q_data.get('q2', ''))
    q_bottom_left = q_data.get('bottom_left', q_data.get('q3', ''))
    q_bottom_right = q_data.get('bottom_right', q_data.get('q4', ''))

    _add_textbox(slide, grid_x + Pt(6), grid_y + Pt(4),
                 Inches(2.5), Inches(0.3),
                 q_top_left, style,
                 font_size=Pt(10),
                 color=style.get('text_light', DEFAULT_STYLE['text_light']))
    _add_textbox(slide, mid_x + Pt(6), grid_y + Pt(4),
                 Inches(2.5), Inches(0.3),
                 q_top_right, style,
                 font_size=Pt(10),
                 color=style.get('text_light', DEFAULT_STYLE['text_light']),
                 bold=True)
    _add_textbox(slide, grid_x + Pt(6), grid_y + grid_h - Inches(0.3),
                 Inches(2.5), Inches(0.3),
                 q_bottom_left, style,
                 font_size=Pt(10),
                 color=style.get('text_light', DEFAULT_STYLE['text_light']))
    _add_textbox(slide, mid_x + Pt(6), grid_y + grid_h - Inches(0.3),
                 Inches(2.5), Inches(0.3),
                 q_bottom_right, style,
                 font_size=Pt(10),
                 color=style.get('text_light', DEFAULT_STYLE['text_light']))

    highlight = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, mid_x, grid_y, grid_w / 2, grid_h / 2
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
    highlight.fill.fore_color.brightness = 0.9
    highlight.line.fill.background()

    for bubble in bubbles:
        bx = float(bubble.get('x', 0.5))
        by = float(bubble.get('y', 0.5))
        bz = float(bubble.get('size', 30))
        b_label = bubble.get('label', '')

        cx = grid_x + int(grid_w / Inches(1) * bx)
        cy = grid_y + int(grid_h / Inches(1) * (1.0 - by))
        bubble_r = Pt(max(bz, 15))

        b_shape = _add_circle(slide, cx, cy, bubble_r, style,
                              fill_color=style.get('secondary', DEFAULT_STYLE['secondary']),
                              line_color=style.get('secondary', DEFAULT_STYLE['secondary']))
        b_shape.fill.fore_color.brightness = 0.3

        _add_textbox(slide, cx - Inches(0.6), cy - Pt(8),
                     Inches(1.2), Inches(0.25),
                     b_label, style,
                     font_size=Pt(8),
                     color=DEFAULT_STYLE['white'],
                     alignment=PP_ALIGN.CENTER)

    if takeaway:
        takeaway_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.35)
        )
        takeaway_shape.fill.solid()
        takeaway_shape.fill.fore_color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
        takeaway_shape.fill.fore_color.brightness = 0.85
        takeaway_shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(6.55), Inches(10.9), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = takeaway
        p.font.size = style.get('font_size_small', Pt(11))
        p.font.color.rgb = style.get('text', DEFAULT_STYLE['text'])
        p.font.name = style.get('font_body', 'Arial')
        p.font.bold = True


# ======================================================================
# 6. Staircase
# ======================================================================

def draw_staircase(slide, data, style):
    """Draw an ascending staircase with steps left to right.

    data = {'steps': [{'label': str, 'items': [str]}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    steps = data.get('steps', [])
    n = len(steps)
    if n == 0:
        return

    start_x = Inches(0.8)
    base_y = Inches(6.2)
    total_w = Inches(11.5)
    step_w = total_w / n - Pt(8)
    max_h = Inches(4.5)

    color_ramp = [
        style.get('secondary', DEFAULT_STYLE['secondary']),
        style.get('accent', DEFAULT_STYLE['accent']),
        style.get('primary', DEFAULT_STYLE['primary']),
    ]

    for i, step in enumerate(steps):
        t = i / max(n - 1, 1)
        h = Inches(1.8) + (max_h - Inches(1.8)) * t
        y = base_y - h
        x = start_x + i * (step_w + Pt(8))

        fill = color_ramp[min(i, len(color_ramp) - 1)]
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, step_w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)

        p = tf.paragraphs[0]
        p.text = step.get('label', '')
        p.font.size = style.get('font_size_body', Pt(13))
        p.font.bold = True
        p.font.color.rgb = DEFAULT_STYLE['white']
        p.font.name = style.get('font_heading', 'Arial')

        for item in step.get('items', []):
            p = tf.add_paragraph()
            p.text = f'  \u2022 {item}'
            p.font.size = style.get('font_size_small', Pt(11))
            p.font.color.rgb = DEFAULT_STYLE['white']
            p.font.name = style.get('font_body', 'Arial')
            p.space_before = Pt(2)

        badge_r = Pt(16)
        badge = _add_circle(slide, x + step_w / 2, y - Pt(10),
                            badge_r, style,
                            fill_color=style.get('accent', DEFAULT_STYLE['accent']),
                            line_color=style.get('accent', DEFAULT_STYLE['accent']))
        _add_textbox(slide, x + step_w / 2 - Pt(10), y - Pt(20),
                     Pt(20), Pt(20),
                     str(i + 1), style,
                     font_size=Pt(12),
                     color=DEFAULT_STYLE['white'],
                     bold=True, alignment=PP_ALIGN.CENTER)

    add_footer_note(slide, '\u9010\u6b65\u5347\u7ea7\uff0c\u6301\u7eed\u8fdb\u6b65', style)


# ======================================================================
# 7. Split-Comparison
# ======================================================================

def draw_split_comparison(slide, data, style):
    """Draw a two-panel comparison with vertical divider.

    data = {'left': {'label': str, 'items': [str]},
            'right': {'label': str, 'items': [str]}}
    """
    add_title_bar(slide, data.get('title', ''), style)

    left = data.get('left', {'label': '', 'items': []})
    right = data.get('right', {'label': '', 'items': []})

    divider_x = Inches(6.666)
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, divider_x - Pt(1.5), Inches(1.2), Pt(3), Inches(5.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = style.get('border', DEFAULT_STYLE['border'])
    divider.line.fill.background()

    left_heading = left.get('label', '')
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.4), Inches(0.4),
                 left_heading, style,
                 font_size=style.get('font_size_heading', Pt(18)),
                 color=style.get('primary', DEFAULT_STYLE['primary']),
                 bold=True, alignment=PP_ALIGN.CENTER)

    left_items = left.get('items', [])
    item_y = Inches(1.75)
    for i, item in enumerate(left_items):
        card_fill = style.get('secondary', DEFAULT_STYLE['secondary'])
        card_fill_light = RGBColor(
            min(card_fill[0] + 100, 255),
            min(card_fill[1] + 100, 255),
            min(card_fill[2] + 100, 255)
        )
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), item_y, Inches(5.4), Inches(0.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_fill_light
        card.line.color.rgb = card_fill
        card.line.width = Pt(0.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        p = tf.paragraphs[0]
        p.text = f'\u2022 {item}'
        p.font.size = style.get('font_size_body', Pt(13))
        p.font.color.rgb = style.get('text', DEFAULT_STYLE['text'])
        p.font.name = style.get('font_body', 'Arial')

        item_y += Inches(0.65)

    right_heading = right.get('label', '')
    _add_textbox(slide, Inches(7.2), Inches(1.2), Inches(5.4), Inches(0.4),
                 right_heading, style,
                 font_size=style.get('font_size_heading', Pt(18)),
                 color=style.get('accent', DEFAULT_STYLE['accent']),
                 bold=True, alignment=PP_ALIGN.CENTER)

    right_items = right.get('items', [])
    item_y = Inches(1.75)
    for i, item in enumerate(right_items):
        card_fill = style.get('accent', DEFAULT_STYLE['accent'])
        card_fill_light = RGBColor(
            min(card_fill[0] + 100, 255),
            min(card_fill[1] + 100, 255),
            min(card_fill[2] + 100, 255)
        )
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.2), item_y, Inches(5.4), Inches(0.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_fill_light
        card.line.color.rgb = card_fill
        card.line.width = Pt(0.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        p = tf.paragraphs[0]
        p.text = f'\u2022 {item}'
        p.font.size = style.get('font_size_body', Pt(13))
        p.font.color.rgb = style.get('text', DEFAULT_STYLE['text'])
        p.font.name = style.get('font_body', 'Arial')

        item_y += Inches(0.65)

    add_footer_note(slide, '\u5de6\u53f3\u5bf9\u6bd4\uff0c\u5dee\u5f02\u7acb\u73b0', style)


# ======================================================================
# 8. Data-Card-Grid
# ======================================================================

def draw_data_card_grid(slide, data, style):
    """Draw a grid of metric cards with large numbers.

    data = {'layout': '2x2'|'1x4', 'cards': [{'label': str, 'value': str, 'unit': str}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    cards = data.get('cards', [])
    layout = data.get('layout', '2x2')
    n = len(cards)
    if n == 0:
        return

    if layout == '1x4':
        cols = min(n, 4)
        rows = 1
    else:
        cols = 2
        rows = min(2, (n + 1) // 2)

    margin_x = Inches(1.0)
    margin_y = Inches(1.5)
    card_w = (Inches(11.3) - (cols - 1) * Pt(12)) / cols
    card_h = (Inches(5.2) - (rows - 1) * Pt(12)) / rows

    for idx, card in enumerate(cards):
        col = idx % cols
        row = idx // cols
        if row >= rows:
            break

        x = margin_x + col * (card_w + Pt(12))
        y = margin_y + row * (card_h + Pt(12))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = style.get('bg', DEFAULT_STYLE['bg'])
        shape.line.color.rgb = style.get('primary', DEFAULT_STYLE['primary'])
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(16)
        tf.margin_right = Pt(16)
        tf.margin_top = Pt(12)

        p = tf.paragraphs[0]
        p.text = str(card.get('value', ''))
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
        p.font.name = style.get('font_heading', 'Arial')

        p = tf.add_paragraph()
        unit_text = card.get('unit', '')
        value_text = str(card.get('value', ''))
        if unit_text:
            p.text = unit_text
            p.font.size = style.get('font_size_small', Pt(11))
            p.font.color.rgb = style.get('text_light', DEFAULT_STYLE['text_light'])
            p.font.name = style.get('font_body', 'Arial')

        p = tf.add_paragraph()
        p.text = card.get('label', '')
        p.font.size = style.get('font_size_body', Pt(13))
        p.font.color.rgb = style.get('text', DEFAULT_STYLE['text'])
        p.font.name = style.get('font_body', 'Arial')
        p.space_before = Pt(6)

    add_footer_note(slide, '\u6570\u636e\u9a71\u52a8\u51b3\u7b56', style)


# ======================================================================
# 9. Layered-Architecture
# ======================================================================

def draw_layered_architecture(slide, data, style):
    """Draw a layered architecture diagram with subcomponents.

    data = {'layers': [{'label': str, 'subcomponents': [{'label': str}, ...]}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    layers = data.get('layers', [])
    n = len(layers)
    if n == 0:
        return

    start_y = Inches(1.2)
    total_h = Inches(5.5)
    layer_h = total_h / n - Pt(6)
    left_w = Inches(2.5)
    right_x = Inches(3.3)
    right_w = Inches(9.0)

    color_ramp = [
        style.get('secondary', DEFAULT_STYLE['secondary']),
        style.get('primary', DEFAULT_STYLE['primary']),
        style.get('primary', DEFAULT_STYLE['primary']),
    ]

    for i, layer in enumerate(layers):
        y = start_y + i * (layer_h + Pt(6))
        fill = color_ramp[min(i, len(color_ramp) - 1)]
        if i > 0:
            brightness = 0.3 + i * 0.3
        else:
            brightness = 0.0

        left_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), y, left_w - Inches(0.1), layer_h
        )
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = fill
        left_shape.line.fill.background()

        tf = left_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        p = tf.paragraphs[0]
        p.text = layer.get('label', '')
        p.font.size = style.get('font_size_small', Pt(11))
        p.font.bold = True
        p.font.color.rgb = DEFAULT_STYLE['white']
        p.font.name = style.get('font_heading', 'Arial')
        p.alignment = PP_ALIGN.CENTER

        subcomponents = layer.get('subcomponents', [])
        sub_n = len(subcomponents)
        if sub_n == 0:
            continue

        sub_w = (right_w - (sub_n - 1) * Pt(8)) / sub_n
        for j, sub in enumerate(subcomponents):
            sx = right_x + j * (sub_w + Pt(8))
            sy = y + Pt(6)
            s_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sub_w, layer_h - Pt(12)
            )
            lighter = RGBColor(
                min(fill[0] + 60, 255),
                min(fill[1] + 60, 255),
                min(fill[2] + 60, 255),
            )
            s_shape.fill.solid()
            s_shape.fill.fore_color.rgb = lighter
            s_shape.line.color.rgb = fill
            s_shape.line.width = Pt(1)

            tf = s_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            p = tf.paragraphs[0]
            p.text = sub.get('label', '')
            p.font.size = style.get('font_size_small', Pt(11))
            p.font.color.rgb = DEFAULT_STYLE['white']
            p.font.name = style.get('font_body', 'Arial')
            p.alignment = PP_ALIGN.CENTER

    add_footer_note(slide, '\u5206\u5c42\u67b6\u6784\uff0c\u6a21\u5757\u89e3\u8026', style)


# ======================================================================
# 10. Filter-Funnel
# ======================================================================

def draw_filter_funnel(slide, data, style):
    """Draw a filter funnel with decreasing-width blocks.

    data = {'layers': [{'label': str, 'items': [str], 'width': float}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    layers = data.get('layers', [])
    n = len(layers)
    if n == 0:
        return

    center_x = Inches(6.666)
    top_y = Inches(1.3)
    total_h = Inches(5.5)
    layer_h = total_h / n - Pt(6)
    max_w = Inches(11.0)

    color_ramp = [
        RGBColor(0xE8, 0xF5, 0xF5),  # light accent
        RGBColor(
            min(style.get('accent', DEFAULT_STYLE['accent'])[0] + 100, 255),
            min(style.get('accent', DEFAULT_STYLE['accent'])[1] + 100, 255),
            min(style.get('accent', DEFAULT_STYLE['accent'])[2] + 100, 255),
        ),
        style.get('accent', DEFAULT_STYLE['accent']),
        style.get('primary', DEFAULT_STYLE['primary']),
    ]

    for i, layer in enumerate(layers):
        w_ratio = layer.get('width', 1.0 - i / max(n, 1))
        w = max_w * w_ratio
        y = top_y + i * (layer_h + Pt(6))

        fill = color_ramp[min(i, len(color_ramp) - 1)]
        text_color = style.get('text', DEFAULT_STYLE['text']) if i == 0 else DEFAULT_STYLE['white']

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            center_x - w / 2, y, w, layer_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        tf.margin_top = Pt(6)

        p = tf.paragraphs[0]
        p.text = layer.get('label', '')
        p.font.size = style.get('font_size_body', Pt(13))
        p.font.bold = True
        p.font.color.rgb = text_color
        p.font.name = style.get('font_heading', 'Arial')

        for item in layer.get('items', []):
            p = tf.add_paragraph()
            p.text = f'  \u2022 {item}'
            p.font.size = style.get('font_size_small', Pt(11))
            p.font.color.rgb = text_color
            p.font.name = style.get('font_body', 'Arial')
            p.space_before = Pt(2)

    add_footer_note(slide, '\u6f0f\u6597\u5f0f\u7b5b\u9009\uff0c\u81ea\u4e0a\u800c\u4e0b\u9010\u7ea7\u7cbe\u7ec6', style)


# ======================================================================
# 11. Overlapping-Spheres
# ======================================================================

def draw_overlapping_spheres(slide, data, style):
    """Draw 2-3 overlapping circles (Venn-style).

    data = {'circles': [{'label': str}, ...],
            'overlap_labels': {...}}
    """
    add_title_bar(slide, data.get('title', ''), style)

    circles = data.get('circles', [])
    overlap_labels = data.get('overlap_labels', {})
    n_circles = len(circles)

    if n_circles == 2:
        c1_x = Inches(4.8)
        c1_y = Inches(3.8)
        c2_x = Inches(8.5)
        c2_y = Inches(3.8)
        r = Inches(1.6)

        fill1 = style.get('secondary', DEFAULT_STYLE['secondary'])
        fill2 = style.get('accent', DEFAULT_STYLE['accent'])

        c1 = _add_circle(slide, c1_x, c1_y, r, style,
                         fill_color=fill1, line_color=fill1)
        c1.fill.fore_color.brightness = 0.3

        c2 = _add_circle(slide, c2_x, c2_y, r, style,
                         fill_color=fill2, line_color=fill2)
        c2.fill.fore_color.brightness = 0.3

        _add_textbox(slide, c1_x - Inches(1.2), c1_y - Pt(10),
                     Inches(2.4), Inches(0.3),
                     circles[0].get('label', '') if n_circles > 0 else '',
                     style, font_size=style.get('font_size_body', Pt(13)),
                     color=style.get('text', DEFAULT_STYLE['text']),
                     bold=True, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, c2_x - Inches(1.2), c2_y - Pt(10),
                     Inches(2.4), Inches(0.3),
                     circles[1].get('label', '') if n_circles > 1 else '',
                     style, font_size=style.get('font_size_body', Pt(13)),
                     color=style.get('text', DEFAULT_STYLE['text']),
                     bold=True, alignment=PP_ALIGN.CENTER)

        overlap_key = '0_1'
        if overlap_key in overlap_labels:
            _add_textbox(slide, Inches(5.8), Inches(3.8 - Pt(8)),
                         Inches(1.6), Inches(0.3),
                         overlap_labels[overlap_key], style,
                         font_size=style.get('font_size_small', Pt(11)),
                         color=style.get('text', DEFAULT_STYLE['text']),
                         bold=True, alignment=PP_ALIGN.CENTER)

    elif n_circles >= 3:
        positions = [
            (Inches(5.5), Inches(3.0)),
            (Inches(3.8), Inches(4.8)),
            (Inches(7.2), Inches(4.8)),
        ]
        fills = [
            style.get('secondary', DEFAULT_STYLE['secondary']),
            style.get('accent', DEFAULT_STYLE['accent']),
            style.get('primary', DEFAULT_STYLE['primary']),
        ]
        r = Inches(1.3)

        for i in range(3):
            if i >= n_circles:
                break
            cx, cy = positions[i]
            fill = fills[min(i, len(fills) - 1)]
            c = _add_circle(slide, cx, cy, r, style,
                            fill_color=fill, line_color=fill)
            c.fill.fore_color.brightness = 0.3

            _add_textbox(slide, cx - Inches(1.0), cy - Pt(10),
                         Inches(2.0), Inches(0.3),
                         circles[i].get('label', ''),
                         style, font_size=style.get('font_size_body', Pt(13)),
                         color=style.get('text', DEFAULT_STYLE['text']),
                         bold=True, alignment=PP_ALIGN.CENTER)

        center_x = Inches(5.5)
        center_y = Inches(4.0)
        overlap_key = '0_1_2'
        if overlap_key in overlap_labels:
            _add_textbox(slide, center_x - Inches(1.0), center_y - Pt(8),
                         Inches(2.0), Inches(0.3),
                         overlap_labels[overlap_key], style,
                         font_size=style.get('font_size_small', Pt(11)),
                         color=style.get('primary', DEFAULT_STYLE['primary']),
                         bold=True, alignment=PP_ALIGN.CENTER)

    add_footer_note(slide, '\u4ea4\u96c6\u4e0e\u5e76\u96c6\u5173\u7cfb', style)


# ======================================================================
# 12. Iterative-Cycle
# ======================================================================

def draw_iterative_cycle(slide, data, style):
    """Draw a circular iterative cycle with directional arrows and center label.

    data = {'steps': [{'label': str}, ...], 'center_label': str}
    """
    add_title_bar(slide, data.get('title', ''), style)

    steps = data.get('steps', [])
    center_label = data.get('center_label', '')
    n = len(steps)
    if n == 0:
        return

    cx = Inches(6.666)
    cy = Inches(4.0)
    radius = Inches(2.3)
    node_r = Inches(0.35)

    positions = []
    for i in range(n):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        nx = cx + radius * math.cos(angle)
        ny = cy + radius * math.sin(angle)
        positions.append((nx, ny))

    accent = style.get('accent', DEFAULT_STYLE['accent'])
    for i in range(n):
        x1, y1 = positions[i]
        x2, y2 = positions[(i + 1) % n]

        dx = float(x2) - float(x1)
        dy = float(y2) - float(y1)
        dist = math.sqrt(dx**2 + dy**2)
        if dist > 0:
            ratio = float(node_r + Pt(5)) / dist
            sx = x1 + dx * ratio
            sy = y1 + dy * ratio
            ex = x2 - dx * ratio
            ey = y2 - dy * ratio

            arrow = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT, sx, sy, ex, ey
            )
            arrow.line.color.rgb = accent
            arrow.line.width = Pt(2)

    for i, step in enumerate(steps):
        nx, ny = positions[i]
        _add_circle(slide, nx, ny, node_r, style,
                    fill_color=style.get('secondary', DEFAULT_STYLE['secondary']),
                    line_color=style.get('secondary', DEFAULT_STYLE['secondary']))

        angle = -math.pi / 2 + 2 * math.pi * i / n
        lx = cx + (radius + Inches(0.7)) * math.cos(angle)
        ly = cy + (radius + Inches(0.7)) * math.sin(angle)

        _add_textbox(slide, lx - Inches(1.0), ly - Pt(10),
                     Inches(2.0), Inches(0.3),
                     step.get('label', ''), style,
                     font_size=style.get('font_size_small', Pt(11)),
                     color=style.get('text', DEFAULT_STYLE['text']),
                     bold=True, alignment=PP_ALIGN.CENTER)

    _add_circle(slide, cx, cy, Inches(0.5), style,
                fill_color=style.get('primary', DEFAULT_STYLE['primary']),
                line_color=style.get('primary', DEFAULT_STYLE['primary']))

    _add_textbox(slide, cx - Inches(1.0), cy - Pt(10),
                 Inches(2.0), Inches(0.3),
                 center_label, style,
                 font_size=style.get('font_size_body', Pt(13)),
                 color=DEFAULT_STYLE['white'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_footer_note(slide, '\u5faa\u73af\u8fed\u4ee3\uff0c\u6301\u7eed\u4f18\u5316', style)


# ======================================================================
# 13. Bridge-and-Gap
# ======================================================================

def draw_bridge_gap(slide, data, style):
    """Draw a bridge-from-current-to-future diagram.

    data = {'current': {'label': str, 'items': [str]},
            'future': {'label': str, 'items': [str]},
            'bridge': {'label': str, 'items': [str]},
            'gap_analysis': str}
    """
    add_title_bar(slide, data.get('title', ''), style)

    current = data.get('current', {'label': '', 'items': []})
    future = data.get('future', {'label': '', 'items': []})
    bridge = data.get('bridge', {'label': '', 'items': []})
    gap_analysis = data.get('gap_analysis', '')

    box_y = Inches(1.4)
    box_h = Inches(4.0)
    box_w = Inches(3.2)

    secondary = style.get('secondary', DEFAULT_STYLE['secondary'])
    add_filled_block(slide, Inches(0.6), box_y, box_w, box_h,
                     current.get('label', ''), current.get('items', []),
                     style, secondary, DEFAULT_STYLE['white'])

    accent = style.get('accent', DEFAULT_STYLE['accent'])
    add_filled_block(slide, Inches(9.5), box_y, box_w, box_h,
                     future.get('label', ''), future.get('items', []),
                     style, accent, DEFAULT_STYLE['white'])

    bridge_y = Inches(2.2)
    bridge_w = Inches(4.8)
    bridge_h = Inches(2.0)

    bridge_card = add_rounded_card(
        slide, Inches(4.2), bridge_y, bridge_w, bridge_h,
        bridge.get('label', ''), bridge.get('items', []),
        style
    )

    arrow_from_x = Inches(4.1)
    arrow_to_x = Inches(9.4)
    arrow_y = Inches(3.2)

    arrow_shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(3.9), arrow_y, Inches(0.4), Inches(0.3)
    )
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = accent
    arrow_shape.line.fill.background()

    arrow_shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(9.1), arrow_y, Inches(0.4), Inches(0.3)
    )
    arrow_shape2.fill.solid()
    arrow_shape2.fill.fore_color.rgb = accent
    arrow_shape2.line.fill.background()

    if gap_analysis:
        gap_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.0), Inches(5.8), Inches(9.3), Inches(0.7)
        )
        gap_shape.fill.solid()
        gap_shape.fill.fore_color.rgb = style.get('bg', DEFAULT_STYLE['bg'])
        gap_shape.line.color.rgb = style.get('border', DEFAULT_STYLE['border'])
        gap_shape.line.width = Pt(1)

        tf = gap_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_top = Pt(4)
        p = tf.paragraphs[0]
        p.text = gap_analysis
        p.font.size = style.get('font_size_small', Pt(11))
        p.font.color.rgb = style.get('text_light', DEFAULT_STYLE['text_light'])
        p.font.name = style.get('font_body', 'Arial')

    add_footer_note(slide, '\u6865\u63a5\u73b0\u72b6\u4e0e\u672a\u6765\uff0c\u8de8\u8d8a\u5dee\u8ddd', style)


# ======================================================================
# 14. Timeline
# ======================================================================

def draw_timeline_diagram(slide, data, style):
    """Draw a horizontal timeline with event nodes.

    data = {'events': [{'date': str, 'desc': str}, ...]}
    """
    add_title_bar(slide, data.get('title', ''), style)

    events = data.get('events', [])
    n = len(events)
    if n == 0:
        return

    line_y = Inches(3.8)
    line_x1 = Inches(1.5)
    line_x2 = Inches(11.8)

    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, line_x1, line_y, line_x2, line_y
    )
    line.line.color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
    line.line.width = Pt(3)

    spacing = (float(line_x2) - float(line_x1)) / max(n - 1, 1)
    node_r = Pt(12)

    for i, event in enumerate(events):
        nx = line_x1 + spacing * i

        _add_circle(slide, nx, line_y, node_r, style,
                    fill_color=style.get('accent', DEFAULT_STYLE['accent']),
                    line_color=style.get('accent', DEFAULT_STYLE['accent']))
        inner_r = Pt(6)
        _add_circle(slide, nx, line_y, inner_r, style,
                    fill_color=style.get('bg', DEFAULT_STYLE['bg']),
                    line_color=style.get('accent', DEFAULT_STYLE['accent']),
                    line_width=Pt(1))

        date_conn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, nx, line_y - node_r,
            nx, line_y - Inches(0.5)
        )
        date_conn.line.color.rgb = style.get('accent', DEFAULT_STYLE['accent'])
        date_conn.line.width = Pt(1)

        _add_textbox(slide, nx - Inches(0.8), line_y - Inches(1.2),
                     Inches(1.6), Inches(0.4),
                     event.get('date', ''), style,
                     font_size=style.get('font_size_small', Pt(11)),
                     color=style.get('accent', DEFAULT_STYLE['accent']),
                     bold=True, alignment=PP_ALIGN.CENTER)

        _add_textbox(slide, nx - Inches(1.4), line_y + Inches(0.4),
                     Inches(2.8), Inches(0.8),
                     event.get('desc', ''), style,
                     font_size=style.get('font_size_body', Pt(13)),
                     color=style.get('text', DEFAULT_STYLE['text']),
                     alignment=PP_ALIGN.CENTER)


# ======================================================================
# Dispatcher
# ======================================================================

DIAGRAM_FUNCTIONS = {
    'pyramid': draw_pyramid,
    'hub_spoke': draw_hub_spoke,
    'dual_gears': draw_dual_gears,
    'tension_triangle': draw_tension_triangle,
    'bubble_matrix': draw_bubble_matrix,
    'staircase': draw_staircase,
    'split_comparison': draw_split_comparison,
    'data_card_grid': draw_data_card_grid,
    'layered_architecture': draw_layered_architecture,
    'filter_funnel': draw_filter_funnel,
    'overlapping_spheres': draw_overlapping_spheres,
    'iterative_cycle': draw_iterative_cycle,
    'bridge_gap': draw_bridge_gap,
    'timeline_diagram': draw_timeline_diagram,
}


def draw_diagram(slide, diagram_type, data=None, style=None):
    """Dispatch to the correct diagram renderer.

    Args:
        slide: python-pptx Slide object.
        diagram_type: str key into DIAGRAM_FUNCTIONS.
        data: dict with type-specific content keys.
        style: optional dict with style overrides.

    Raises:
        KeyError: if diagram_type is not registered.
    """
    func = DIAGRAM_FUNCTIONS.get(diagram_type)
    if func is None:
        raise KeyError(
            f"Unknown diagram type '{diagram_type}'. "
            f"Available: {', '.join(sorted(DIAGRAM_FUNCTIONS.keys()))}"
        )
    resolved_style = _resolve_style(style)
    func(slide, data or {}, resolved_style)
